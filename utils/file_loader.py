import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd


def load_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])


def load_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])


def load_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def load_text_from_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)


def load_text(file_path):
    if file_path.endswith(".pdf"):
        return load_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return load_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return load_text_from_pptx(file_path)
    elif file_path.endswith(".csv"):
        return load_text_from_csv(file_path)
    elif file_path.endswith(".txt") or file_path.endswith(".md"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        return "Unsupported format."